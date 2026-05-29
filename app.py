"""
SlideCraft Lite — minimal slide editor with per-session isolation.

Features:
  • Upload PPTX (single + bulk)
  • Remove NotebookLM logo
  • Region erase (click-drag a box to clone-stamp it out)
  • Reorder / Duplicate / Delete slides
  • Watermark (text + image)
  • Undo / Redo across destructive ops
  • Save snapshot, Reset to originals
  • Present (client-side), Help (client-side)
  • Export (PPTX, PDF) + per-slide PNG
  • Bulk batch — 10 files / 300 MB, async with progress
  • Recent decks (per session)
  • Per-session storage isolation (cookie-scoped)
"""
import os
import sys
import json
import re
import shutil
import subprocess
import threading
import time
import uuid
import datetime
import io
import zipfile
import tempfile
from pathlib import Path
from types import SimpleNamespace

from flask import (Flask, render_template, jsonify, request, send_file,
                   session, Response, abort, send_from_directory)
from pptx import Presentation
from pptx.util import Inches
from werkzeug.utils import secure_filename
from werkzeug.exceptions import RequestEntityTooLarge
from PIL import Image, ImageDraw, ImageFont

Image.MAX_IMAGE_PIXELS = 50_000_000

_SAFE_NAME_RE = re.compile(r"^[A-Za-z0-9_\- ]{1,64}$")
_SID_RE = re.compile(r"^[0-9a-f]{32}$")


def _safe_name(s):
    if not isinstance(s, str):
        return ""
    s = s.strip()
    if not s or not _SAFE_NAME_RE.match(s):
        return ""
    return s if s not in (".", "..") else ""


def _ensure_dict(payload):
    return payload if isinstance(payload, dict) else {}


app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY') or os.urandom(32).hex()
app.permanent_session_lifetime = datetime.timedelta(days=7)
app.config['MAX_CONTENT_LENGTH'] = int(os.environ.get('MAX_UPLOAD_MB', 300)) * 1024 * 1024
BULK_FILE_CAP = 10
BULK_BYTES_CAP = 300 * 1024 * 1024
SESSION_TTL_HOURS = int(os.environ.get('SESSION_TTL_HOURS', 24))
RECENT_KEEP = 5


@app.errorhandler(RequestEntityTooLarge)
def _too_large(e):
    cap_mb = app.config['MAX_CONTENT_LENGTH'] // (1024 * 1024)
    return jsonify({
        "error": f"Upload too large — combined size exceeds {cap_mb} MB cap.",
        "cap_mb": cap_mb,
    }), 413


BASE_DIR = Path(__file__).parent
SESSIONS_DIR = BASE_DIR / "sessions"
SESSIONS_DIR.mkdir(parents=True, exist_ok=True)

SLIDE_W_PX, SLIDE_H_PX = 2134, 1200
_data_lock = threading.RLock()
IS_WINDOWS = sys.platform == "win32"
IS_MACOS = sys.platform == "darwin"


# ── Session paths ────────────────────────────────────────────────────────
def _sid():
    sid = session.get('sid')
    if not sid or not _SID_RE.match(sid):
        sid = uuid.uuid4().hex
        session['sid'] = sid
        session.permanent = True
    return sid


def P():
    """Return a SimpleNamespace of paths for the current session."""
    root = SESSIONS_DIR / _sid()
    p = SimpleNamespace(
        root=root,
        slides=root / "slides",
        originals=root / "slides" / "_originals",
        uploads=root / "uploads",
        exports=root / "exports",
        history=root / "history",
        recent=root / "recent",
        log=root / "ops_log.json",
        deck_info=root / "current_deck.txt",
        recent_index=root / "recent.json",
        last_saved=root / "last_saved.txt",
        jobs=root / "jobs",
    )
    for d in (p.root, p.slides, p.originals, p.uploads, p.exports,
              p.history, p.recent, p.jobs):
        d.mkdir(parents=True, exist_ok=True)
    return p


def _mark_saved():
    try:
        P().last_saved.write_text(datetime.datetime.now().isoformat(timespec='seconds'))
    except OSError:
        pass


def _set_deck_name(name):
    try:
        P().deck_info.write_text(name)
    except OSError:
        pass


def _get_deck_name():
    try:
        f = P().deck_info
        return f.read_text().strip() if f.exists() else ""
    except OSError:
        return ""


def _get_slide_files():
    return sorted(P().slides.glob("slide-*.jpg"))


# ── NotebookLM logo eraser ────────────────────────────────────────────────
LOGO_REF_W, LOGO_REF_H = 1376, 768
LOGO_W, LOGO_H = 145, 28


def _erase_logo(img):
    w, h = img.size
    logo_w = int(LOGO_W * w / LOGO_REF_W)
    logo_h = int(LOGO_H * h / LOGO_REF_H)
    src_y = max(0, h - logo_h - 1)
    strip = img.crop((w - logo_w, src_y, w, src_y + 1))
    patch = strip.resize((logo_w, logo_h), Image.NEAREST)
    img.paste(patch, (w - logo_w, h - logo_h))
    return img


def remove_logos_batch(slide_files):
    for p in slide_files:
        img = Image.open(p).convert("RGB")
        img = _erase_logo(img)
        img.save(str(p), quality=95)


def _erase_region(img, box):
    """Erase an arbitrary box from img by tiling the row immediately above
    (or below if at the very top) down across the region. Same trick as
    the logo eraser, but with a caller-supplied box."""
    x, y, bw, bh = box
    w, h = img.size
    x = max(0, min(w - 1, int(x)))
    y = max(0, min(h - 1, int(y)))
    bw = max(1, min(w - x, int(bw)))
    bh = max(1, min(h - y, int(bh)))
    src_y = y - 1 if y > 0 else min(h - 1, y + bh)
    strip = img.crop((x, src_y, x + bw, src_y + 1))
    patch = strip.resize((bw, bh), Image.NEAREST)
    img.paste(patch, (x, y))
    return img


# ── PPTX → slide image conversion ────────────────────────────────────────
def _find_libreoffice():
    candidates = ["libreoffice", "soffice"]
    if IS_WINDOWS:
        candidates.extend([
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            str(Path.home() / "AppData" / "Local" / "Programs" / "LibreOffice" / "program" / "soffice.exe"),
        ])
    elif IS_MACOS:
        candidates.append("/Applications/LibreOffice.app/Contents/MacOS/soffice")
    else:
        candidates.extend(["/usr/bin/soffice", "/snap/bin/libreoffice", "/usr/bin/libreoffice"])
    for c in candidates:
        if shutil.which(c) or Path(c).exists():
            return c
    return None


def _convert_pptx_to_images_libreoffice(pptx_path, output_dir):
    lo = _find_libreoffice()
    if not lo:
        raise RuntimeError("LibreOffice not found")
    tmp_dir = Path(tempfile.mkdtemp())
    try:
        subprocess.run(
            [lo, "--headless", "--convert-to", "pdf", "--outdir", str(tmp_dir), str(pptx_path)],
            check=True, timeout=120,
        )
        pdfs = list(tmp_dir.glob("*.pdf"))
        if not pdfs:
            raise RuntimeError("PDF conversion produced no output")
        try:
            from pdf2image import convert_from_path
            for i, img in enumerate(convert_from_path(str(pdfs[0]), dpi=200)):
                img.convert("RGB").save(str(output_dir / f"slide-{i+1:02d}.jpg"), "JPEG", quality=95)
        except ImportError:
            import fitz
            doc = fitz.open(str(pdfs[0]))
            mat = fitz.Matrix(2.0, 2.0)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat)
                (output_dir / f"slide-{i+1:02d}.jpg").write_bytes(pix.tobytes("jpeg"))
            doc.close()
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def _convert_pptx_to_images_pillow(pptx_path, output_dir):
    prs = Presentation(str(pptx_path))
    sw_emu, sh_emu = prs.slide_width, prs.slide_height
    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (SLIDE_W_PX, SLIDE_H_PX), (255, 255, 255))
        for shape in slide.shapes:
            if shape.shape_type == 13:
                shape_img = Image.open(io.BytesIO(shape.image.blob))
                left = int(shape.left / sw_emu * SLIDE_W_PX) if sw_emu else 0
                top = int(shape.top / sh_emu * SLIDE_H_PX) if sh_emu else 0
                sw = int(shape.width / sw_emu * SLIDE_W_PX) if sw_emu else SLIDE_W_PX
                sh = int(shape.height / sh_emu * SLIDE_H_PX) if sh_emu else SLIDE_H_PX
                img.paste(shape_img.resize((sw, sh), Image.BILINEAR), (left, top))
        img.save(str(output_dir / f"slide-{i+1:02d}.jpg"), "JPEG", quality=95)


def process_uploaded_pptx(pptx_path):
    p = P()
    p.slides.mkdir(parents=True, exist_ok=True)
    stage_dir = Path(tempfile.mkdtemp(prefix="slides_stage_"))
    try:
        try:
            _convert_pptx_to_images_libreoffice(pptx_path, stage_dir)
        except (RuntimeError, FileNotFoundError, subprocess.SubprocessError, OSError):
            _convert_pptx_to_images_pillow(pptx_path, stage_dir)
        staged = sorted(stage_dir.glob("slide-*.jpg"))
        if not staged:
            raise RuntimeError("Conversion produced no slide images")
        for f in p.slides.glob("slide-*.jpg"):
            f.unlink()
        for f in p.originals.glob("slide-*.jpg"):
            f.unlink()
        for sf in staged:
            shutil.copy2(str(sf), str(p.originals / sf.name))
            shutil.move(str(sf), str(p.slides / sf.name))
    finally:
        shutil.rmtree(stage_dir, ignore_errors=True)


def _rebuild_pptx_from_images(slide_images, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for sf in slide_images:
        slide = prs.slides.add_slide(blank)
        pic = slide.shapes.add_picture(str(sf), 0, 0, prs.slide_width, prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    prs.save(str(output_path))


# ── Op log + history snapshot system ──────────────────────────────────────
def _load_log():
    with _data_lock:
        log = P().log
        if log.exists():
            try:
                return json.loads(log.read_text())
            except (json.JSONDecodeError, OSError):
                return []
        return []


def _save_log(entries):
    with _data_lock:
        P().log.write_text(json.dumps(entries, indent=2))


def _append_log(entry):
    with _data_lock:
        entries = _load_log()
        entries = [e for e in entries if not e.get("undone")]
        entries.append(entry)
        if len(entries) > 100:
            entries = entries[-100:]
        P().log.write_text(json.dumps(entries, indent=2))


def _snapshot(reason="op"):
    try:
        p = P()
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S") + f"_{uuid.uuid4().hex[:4]}"
        version_dir = p.history / ts
        version_dir.mkdir(exist_ok=True)
        for sf in _get_slide_files():
            shutil.copy2(str(sf), str(version_dir / sf.name))
        (version_dir / "_reason.txt").write_text(reason)
        return ts
    except OSError:
        return None


def _restore_snapshot(version_dir):
    snap_slides = sorted(version_dir.glob("slide-*.jpg"))
    if not snap_slides:
        raise RuntimeError("Snapshot contains no slide images")
    slides_dir = P().slides
    for f in slides_dir.glob("slide-*.jpg"):
        try:
            f.unlink()
        except OSError:
            pass
    for sf in snap_slides:
        try:
            shutil.copy2(str(sf), str(slides_dir / sf.name))
        except OSError as e:
            raise RuntimeError(f"Copy failed for {sf.name}: {e}") from e


def _log_op(kind, *, text, snapshot, count=1):
    if not snapshot:
        return None
    entry_id = uuid.uuid4().hex[:12]
    _append_log({
        "id": entry_id,
        "kind": kind,
        "text": text,
        "count": count,
        "snapshot": snapshot,
        "timestamp": datetime.datetime.now().isoformat(timespec="seconds"),
    })
    _mark_saved()
    return entry_id


# ── Recent decks (per session) ────────────────────────────────────────────
def _load_recent():
    f = P().recent_index
    if not f.exists():
        return []
    try:
        return json.loads(f.read_text())
    except (json.JSONDecodeError, OSError):
        return []


def _save_recent(items):
    P().recent_index.write_text(json.dumps(items, indent=2))


def _push_recent(orig_name, stored_path):
    items = _load_recent()
    items = [it for it in items if it.get("file") != Path(stored_path).name]
    items.insert(0, {
        "name": orig_name,
        "file": Path(stored_path).name,
        "uploaded": datetime.datetime.now().isoformat(timespec="seconds"),
    })
    while len(items) > RECENT_KEEP:
        dropped = items.pop()
        try:
            (P().recent / dropped["file"]).unlink(missing_ok=True)
        except OSError:
            pass
    _save_recent(items)


# ── Routes ────────────────────────────────────────────────────────────────
@app.route("/")
def index():
    _sid()
    slide_files = _get_slide_files()
    slides = [{"index": i + 1, "file": f.name} for i, f in enumerate(slide_files)]
    return render_template("index.html", slides=slides, num_slides=len(slide_files))


@app.route("/sl/<path:name>")
def serve_slide(name):
    if not re.fullmatch(r"slide-\d{2,4}\.jpg", name):
        abort(404)
    f = P().slides / name
    if not f.exists():
        abort(404)
    return send_from_directory(str(P().slides), name, max_age=0)


@app.route("/api/upload", methods=["POST"])
def upload_pptx():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    f = request.files["file"]
    if not f.filename.lower().endswith(".pptx"):
        return jsonify({"error": "Only .pptx files are supported"}), 400
    safe_name = secure_filename(f.filename)
    p = P()
    save_path = p.uploads / safe_name
    f.save(str(save_path))
    try:
        process_uploaded_pptx(save_path)
    except Exception as e:
        return jsonify({"error": f"Processing failed: {e}"}), 500
    _set_deck_name(f.filename or safe_name)
    # Stash a copy into recent/
    recent_copy = p.recent / f"{uuid.uuid4().hex[:8]}_{safe_name}"
    try:
        shutil.copy2(str(save_path), str(recent_copy))
        _push_recent(f.filename or safe_name, recent_copy)
    except OSError:
        pass
    _mark_saved()
    return jsonify({"ok": True, "num_slides": len(_get_slide_files()),
                    "deck_name": _get_deck_name()})


@app.route("/api/remove-logo", methods=["POST"])
def remove_logo_from_existing():
    slides = sorted(P().slides.glob("slide-*.jpg"))
    if not slides:
        return jsonify({"ok": True, "count": 0})
    snapshot = _snapshot("remove-logo")
    remove_logos_batch(slides)
    log_id = _log_op("remove-logo",
                     text=f"Removed NotebookLM logo from {len(slides)} slide(s)",
                     snapshot=snapshot, count=len(slides))
    return jsonify({"ok": True, "count": len(slides),
                    "snapshot": snapshot, "log_id": log_id})


@app.route("/api/region-erase", methods=["POST"])
def region_erase():
    """Erase an arbitrary box from a slide using the existing tile trick.
    Body: { slide_num, x, y, w, h }  — all normalised (0..1) of the slide.
    """
    payload = _ensure_dict(request.json)
    slide_files = _get_slide_files()
    try:
        n = int(payload.get("slide_num", 0))
    except (TypeError, ValueError):
        return jsonify({"error": "Invalid slide_num"}), 400
    if n < 1 or n > len(slide_files):
        return jsonify({"error": "Invalid slide_num"}), 400

    def _f(k, lo=0.0, hi=1.0):
        try:
            return max(lo, min(hi, float(payload.get(k))))
        except (TypeError, ValueError):
            return None
    x, y, bw, bh = _f("x"), _f("y"), _f("w"), _f("h")
    if None in (x, y, bw, bh):
        return jsonify({"error": "x,y,w,h required"}), 400
    if bw < 0.005 or bh < 0.005:
        return jsonify({"error": "Selection too small"}), 400

    sf = slide_files[n - 1]
    snapshot = _snapshot("region-erase")
    img = Image.open(sf).convert("RGB")
    iw, ih = img.size
    box = (x * iw, y * ih, bw * iw, bh * ih)
    img = _erase_region(img, box)
    img.save(str(sf), quality=95)
    log_id = _log_op("region-erase",
                     text=f"Erased region on slide {n}",
                     snapshot=snapshot)
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id})


@app.route("/api/save", methods=["POST"])
def save_snapshot():
    if not _get_slide_files():
        return jsonify({"error": "Nothing to save — upload a PPTX first"}), 400
    snapshot = _snapshot("save")
    log_id = _log_op("save", text="Manual save",
                     snapshot=snapshot, count=len(_get_slide_files()))
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id})


@app.route("/api/reset-all", methods=["POST"])
def reset_all():
    p = P()
    originals = sorted(p.originals.glob("slide-*.jpg"))
    source_label = "originals"
    if not originals:
        snapshot_dirs = sorted(
            [d for d in p.history.iterdir() if d.is_dir()],
            key=lambda d: d.stat().st_mtime
        ) if p.history.exists() else []
        if not snapshot_dirs:
            return jsonify({"error": "No originals or history snapshots — re-upload the PPTX."}), 400
        originals = sorted(snapshot_dirs[0].glob("slide-*.jpg"))
        if not originals:
            return jsonify({"error": "Oldest snapshot is empty"}), 400
        source_label = f"oldest snapshot ({snapshot_dirs[0].name})"

    snapshot = _snapshot("reset-all")
    for f in p.slides.glob("slide-*.jpg"):
        f.unlink()
    for o in originals:
        shutil.copy2(str(o), str(p.slides / o.name))
    if p.log.exists():
        p.log.unlink()
    _mark_saved()
    return jsonify({"ok": True, "slides_restored": len(originals),
                    "snapshot": snapshot, "source": source_label})


# ── Export ────────────────────────────────────────────────────────────────
@app.route("/api/export", methods=["POST"])
def export_pptx():
    slide_files = _get_slide_files()
    if not slide_files:
        return jsonify({"error": "No slides to export"}), 400
    out_path = P().exports / f"SlideCraft_Export_{uuid.uuid4().hex[:8]}.pptx"
    _rebuild_pptx_from_images(slide_files, out_path)
    return send_file(str(out_path), as_attachment=True, download_name="SlideCraft_Export.pptx")


@app.route("/api/export-pdf", methods=["POST"])
def export_pdf():
    slide_files = _get_slide_files()
    if not slide_files:
        return jsonify({"error": "No slides to export"}), 400
    first_img = Image.open(slide_files[0]).convert("RGB")
    rest = (Image.open(sf).convert("RGB") for sf in slide_files[1:])
    out_path = P().exports / f"Slides_Export_{uuid.uuid4().hex[:8]}.pdf"
    first_img.save(str(out_path), save_all=True, append_images=rest, resolution=150)
    return send_file(str(out_path), as_attachment=True, download_name="Slides_Export.pdf")


# ── Bulk batch (async with progress) ─────────────────────────────────────
_jobs_lock = threading.RLock()


def _job_path(sid, job_id):
    return SESSIONS_DIR / sid / "jobs" / f"{job_id}.json"


def _job_write(sid, job_id, **fields):
    with _jobs_lock:
        f = _job_path(sid, job_id)
        f.parent.mkdir(parents=True, exist_ok=True)
        cur = {}
        if f.exists():
            try:
                cur = json.loads(f.read_text())
            except (json.JSONDecodeError, OSError):
                cur = {}
        cur.update(fields)
        f.write_text(json.dumps(cur))


def _job_read(sid, job_id):
    f = _job_path(sid, job_id)
    if not f.exists():
        return None
    try:
        return json.loads(f.read_text())
    except (json.JSONDecodeError, OSError):
        return None


def _bulk_worker(sid, job_id, file_data, exports_dir):
    tmp_dir = Path(tempfile.mkdtemp())
    cleaned_paths = []
    results = []
    total = len(file_data)
    try:
        _job_write(sid, job_id, status="running", total=total, done=0,
                   message="Starting...")
        for i, (fname, content) in enumerate(file_data, 1):
            _job_write(sid, job_id, message=f"Processing {fname} ({i}/{total})",
                       done=i - 1)
            try:
                input_path = tmp_dir / fname
                input_path.write_bytes(content)
                file_slides_dir = tmp_dir / f"slides_{fname}"
                file_slides_dir.mkdir(exist_ok=True)
                try:
                    _convert_pptx_to_images_libreoffice(input_path, file_slides_dir)
                except (RuntimeError, FileNotFoundError, subprocess.SubprocessError, OSError):
                    _convert_pptx_to_images_pillow(input_path, file_slides_dir)
                slide_images = sorted(file_slides_dir.glob("slide-*.jpg"))
                remove_logos_batch(slide_images)
                output_path = tmp_dir / f"clean_{fname}"
                _rebuild_pptx_from_images(slide_images, output_path)
                cleaned_paths.append((output_path, f"clean_{fname}"))
                results.append({"file": fname, "status": "ok", "slides": len(slide_images)})
            except Exception as e:
                results.append({"file": fname, "status": "error", "error": str(e)})
            _job_write(sid, job_id, done=i, results=results)

        if not cleaned_paths:
            _job_write(sid, job_id, status="error",
                       message="No files were processed", results=results)
            return

        zip_name = f"bulk_cleaned_{uuid.uuid4().hex[:8]}.zip"
        zip_path = exports_dir / zip_name
        with zipfile.ZipFile(str(zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
            for fpath, arcname in cleaned_paths:
                zf.write(str(fpath), arcname)
        _job_write(sid, job_id, status="done", message="Complete",
                   download=zip_name, results=results)
    except Exception as e:
        _job_write(sid, job_id, status="error", message=str(e))
    finally:
        shutil.rmtree(str(tmp_dir), ignore_errors=True)


@app.route("/api/batch/remove-logo", methods=["POST"])
def batch_remove_logo():
    """Kicks off a background job and returns a job_id immediately.
    Client polls /api/job/<id> for progress, then GETs /api/job/<id>/download.
    """
    files = request.files.getlist("files")
    if not files or not any(f.filename for f in files):
        return jsonify({"error": "No files uploaded"}), 400
    pptx_files = [f for f in files if f.filename and f.filename.lower().endswith(".pptx")]
    if not pptx_files:
        return jsonify({"error": "No .pptx files found"}), 400
    if len(pptx_files) > BULK_FILE_CAP:
        return jsonify({"error": f"Maximum {BULK_FILE_CAP} files allowed"}), 400

    total_bytes = 0
    file_data = []
    for f in pptx_files:
        fname = secure_filename(f.filename)
        if not fname:
            continue
        content = f.read()
        total_bytes += len(content)
        file_data.append((fname, content))

    if total_bytes > BULK_BYTES_CAP:
        cap_mb = BULK_BYTES_CAP // (1024 * 1024)
        total_mb = total_bytes // (1024 * 1024)
        return jsonify({
            "error": f"Combined bulk size {total_mb} MB exceeds {cap_mb} MB cap.",
        }), 413

    sid = _sid()
    p = P()
    job_id = uuid.uuid4().hex[:12]
    _job_write(sid, job_id, status="queued", total=len(file_data),
               done=0, message="Queued")
    th = threading.Thread(target=_bulk_worker,
                          args=(sid, job_id, file_data, p.exports), daemon=True)
    th.start()
    return jsonify({"ok": True, "job_id": job_id})


@app.route("/api/job/<job_id>", methods=["GET"])
def job_status(job_id):
    if not re.fullmatch(r"[0-9a-f]{1,32}", job_id):
        return jsonify({"error": "Invalid job_id"}), 400
    sid = _sid()
    state = _job_read(sid, job_id)
    if state is None:
        return jsonify({"error": "Unknown job"}), 404
    return jsonify(state)


@app.route("/api/job/<job_id>/download", methods=["GET"])
def job_download(job_id):
    if not re.fullmatch(r"[0-9a-f]{1,32}", job_id):
        abort(400)
    sid = _sid()
    state = _job_read(sid, job_id)
    if not state or state.get("status") != "done":
        abort(404)
    fname = state.get("download")
    if not fname:
        abort(404)
    f = P().exports / fname
    if not f.exists():
        abort(404)
    return send_file(str(f), as_attachment=True,
                     download_name="SlideCraft_Bulk_Cleaned.zip")


# ── Undo / Redo ──────────────────────────────────────────────────────────
@app.route("/api/ops/state", methods=["GET"])
def ops_state():
    entries = _load_log()
    last_active = None
    last_undone = None
    for e in entries:
        if e.get("undone"):
            last_undone = e
        else:
            last_active = e
    last_saved = ""
    try:
        f = P().last_saved
        if f.exists():
            last_saved = f.read_text().strip()
    except OSError:
        pass
    return jsonify({
        "can_undo": last_active is not None,
        "can_redo": last_undone is not None,
        "undo_label": (last_active and last_active.get("text")) or None,
        "redo_label": (last_undone and last_undone.get("text")) or None,
        "last_saved": last_saved,
    })


@app.route("/api/ops/undo", methods=["POST"])
def ops_undo():
    with _data_lock:
        p = P()
        entries = _load_log()
        target_idx = None
        for i in range(len(entries) - 1, -1, -1):
            if not entries[i].get("undone"):
                target_idx = i
                break
        if target_idx is None:
            return jsonify({"ok": False, "reason": "Nothing to undo"})
        target = entries[target_idx]
        snapshot_id = target.get("snapshot")
        if not snapshot_id:
            return jsonify({"error": "Entry has no snapshot"}), 400
        version_dir = (p.history / snapshot_id).resolve()
        if p.history.resolve() not in version_dir.parents:
            return jsonify({"error": "Invalid snapshot path"}), 400
        if not version_dir.exists():
            return jsonify({"error": "Snapshot no longer exists"}), 404
        try:
            redo_snapshot = _snapshot("redo-snap")
            _restore_snapshot(version_dir)
        except (OSError, RuntimeError) as e:
            return jsonify({"error": f"Restore failed: {e}"}), 500
        entries[target_idx]["undone"] = True
        entries[target_idx]["redo_snapshot"] = redo_snapshot
        p.log.write_text(json.dumps(entries, indent=2))
    _mark_saved()
    return jsonify({"ok": True, "kind": target.get("kind"),
                    "text": target.get("text"), "id": target["id"]})


@app.route("/api/ops/redo", methods=["POST"])
def ops_redo():
    with _data_lock:
        p = P()
        entries = _load_log()
        target_idx = None
        for i in range(len(entries) - 1, -1, -1):
            if entries[i].get("undone"):
                target_idx = i
                break
        if target_idx is None:
            return jsonify({"ok": False, "reason": "Nothing to redo"})
        target = entries[target_idx]
        redo_snap = target.get("redo_snapshot")
        if not redo_snap:
            return jsonify({"error": "Undone entry has no redo snapshot"}), 400
        version_dir = (p.history / redo_snap).resolve()
        if p.history.resolve() not in version_dir.parents:
            return jsonify({"error": "Invalid redo snapshot path"}), 400
        if not version_dir.exists():
            return jsonify({"error": "Redo snapshot no longer exists"}), 404
        try:
            _restore_snapshot(version_dir)
        except (OSError, RuntimeError) as e:
            return jsonify({"error": f"Restore failed: {e}"}), 500
        entries[target_idx]["undone"] = False
        entries[target_idx].pop("redo_snapshot", None)
        p.log.write_text(json.dumps(entries, indent=2))
    _mark_saved()
    return jsonify({"ok": True, "kind": target.get("kind"),
                    "text": target.get("text"), "id": target["id"]})


# ── Deck info ─────────────────────────────────────────────────────────────
@app.route("/api/deck/info", methods=["GET"])
def deck_info():
    return jsonify({
        "name": _get_deck_name(),
        "num_slides": len(_get_slide_files()),
    })


# ── Recent decks ──────────────────────────────────────────────────────────
@app.route("/api/recent", methods=["GET"])
def recent_list():
    return jsonify({"items": _load_recent()})


@app.route("/api/recent/load", methods=["POST"])
def recent_load():
    payload = _ensure_dict(request.json)
    fname = payload.get("file", "")
    if not re.fullmatch(r"[A-Za-z0-9_.\-]{1,128}", fname or ""):
        return jsonify({"error": "Invalid filename"}), 400
    f = P().recent / fname
    if not f.exists():
        return jsonify({"error": "File not found"}), 404
    items = _load_recent()
    orig = next((it for it in items if it.get("file") == fname), None)
    try:
        process_uploaded_pptx(f)
    except Exception as e:
        return jsonify({"error": f"Processing failed: {e}"}), 500
    if orig:
        _set_deck_name(orig.get("name", fname))
        _push_recent(orig.get("name", fname), f)
    _mark_saved()
    return jsonify({"ok": True, "num_slides": len(_get_slide_files()),
                    "deck_name": _get_deck_name()})


@app.route("/api/recent/delete", methods=["POST"])
def recent_delete():
    payload = _ensure_dict(request.json)
    fname = payload.get("file", "")
    if not re.fullmatch(r"[A-Za-z0-9_.\-]{1,128}", fname or ""):
        return jsonify({"error": "Invalid filename"}), 400
    items = [it for it in _load_recent() if it.get("file") != fname]
    _save_recent(items)
    try:
        (P().recent / fname).unlink(missing_ok=True)
    except OSError:
        pass
    return jsonify({"ok": True, "items": items})


# ── Per-slide ops: delete, duplicate, download as PNG ────────────────────
def _renumber_slides():
    slides_dir = P().slides
    files = sorted(slides_dir.glob("slide-*.jpg"))
    tmp_paths = []
    for f in files:
        tmp = f.with_name(f.name + ".tmp")
        f.rename(tmp)
        tmp_paths.append(tmp)
    for i, tmp in enumerate(tmp_paths, 1):
        tmp.rename(slides_dir / f"slide-{i:02d}.jpg")


@app.route("/api/slide/<int:num>/delete", methods=["POST"])
def delete_slide(num):
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400
    if len(slide_files) <= 1:
        return jsonify({"error": "Can't delete the last remaining slide"}), 400
    snapshot = _snapshot("delete")
    slide_files[num - 1].unlink()
    _renumber_slides()
    log_id = _log_op("delete", text=f"Deleted slide {num}", snapshot=snapshot)
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id,
                    "num_slides": len(_get_slide_files())})


@app.route("/api/slide/<int:num>/duplicate", methods=["POST"])
def duplicate_slide(num):
    slides_dir = P().slides
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400
    snapshot = _snapshot("duplicate")
    src = slide_files[num - 1]
    for i in range(len(slide_files), num, -1):
        old = slides_dir / f"slide-{i:02d}.jpg"
        new = slides_dir / f"slide-{i+1:02d}.jpg"
        old.rename(new)
    shutil.copy2(str(src), str(slides_dir / f"slide-{num+1:02d}.jpg"))
    log_id = _log_op("duplicate", text=f"Duplicated slide {num}", snapshot=snapshot)
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id,
                    "num_slides": len(_get_slide_files())})


@app.route("/api/slide/<int:num>/download.png", methods=["GET"])
def download_slide_png(num):
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 404
    img = Image.open(slide_files[num - 1]).convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return send_file(buf, mimetype="image/png", as_attachment=True,
                     download_name=f"slide-{num:02d}.png")


# ── Reorder slides via drag-drop ─────────────────────────────────────────
@app.route("/api/reorder", methods=["POST"])
def reorder_slides():
    payload = _ensure_dict(request.json)
    new_order = payload.get("order", [])
    slides_dir = P().slides
    slide_files = _get_slide_files()
    if sorted(new_order) != list(range(1, len(slide_files) + 1)):
        return jsonify({"error": "Invalid order"}), 400
    snapshot = _snapshot("reorder")
    tmp_dir = slides_dir / "_reorder_tmp"
    try:
        if tmp_dir.exists():
            shutil.rmtree(str(tmp_dir))
        tmp_dir.mkdir()
        for new_idx, old_idx in enumerate(new_order, 1):
            src = slides_dir / f"slide-{old_idx:02d}.jpg"
            dst = tmp_dir / f"slide-{new_idx:02d}.jpg"
            if not src.exists():
                raise RuntimeError(f"Missing slide-{old_idx:02d}.jpg")
            shutil.copy2(str(src), str(dst))
        for f in slides_dir.glob("slide-*.jpg"):
            f.unlink()
        for f in tmp_dir.glob("slide-*.jpg"):
            shutil.move(str(f), str(slides_dir / f.name))
    except Exception as e:
        return jsonify({"error": f"Reorder failed: {e}"}), 500
    finally:
        shutil.rmtree(str(tmp_dir), ignore_errors=True)
    log_id = _log_op("reorder", text=f"Reordered {len(new_order)} slides",
                     snapshot=snapshot, count=len(new_order))
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id})


# ── Watermark insertion (text + image) ───────────────────────────────────
def _parse_hex(s, default=(128, 128, 128)):
    try:
        s = s.lstrip("#")
        if len(s) < 6:
            return default
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except (ValueError, TypeError, AttributeError):
        return default


def _load_bold_font(size):
    candidates = []
    if IS_WINDOWS:
        win_fonts = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
        candidates.extend([win_fonts / n for n in
                           ("arialbd.ttf", "segoeuib.ttf", "calibrib.ttf", "verdanab.ttf")])
    elif IS_MACOS:
        candidates.extend(Path(d) / n for d in
                          ("/System/Library/Fonts", "/System/Library/Fonts/Supplemental", "/Library/Fonts")
                          for n in ("Helvetica.ttc", "Arial Bold.ttf", "Arial.ttf"))
    else:
        for d in ("/usr/share/fonts", "/usr/local/share/fonts"):
            base = Path(d)
            if base.exists():
                for f in base.rglob("DejaVuSans-Bold.ttf"):
                    candidates.append(f); break
                for f in base.rglob("LiberationSans-Bold.ttf"):
                    candidates.append(f); break
    for c in candidates:
        if c.exists():
            try:
                return ImageFont.truetype(str(c), size)
            except (OSError, IOError):
                continue
    return ImageFont.load_default()


def _draw_text_watermark(img, *, text, color, opacity, position, font_scale, rotation,
                          tile_spacing, custom_x=0.5, custom_y=0.5):
    w, h = img.size
    alpha = max(0, min(255, int(opacity * 255)))
    r, g, b = color
    fill = (r, g, b, alpha)
    font_size = max(8, int(w * font_scale))
    font = _load_bold_font(font_size)
    bbox = font.getbbox(text)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))

    def _stamp(x, y):
        if rotation % 360 == 0:
            ImageDraw.Draw(layer).text((x, y), text, fill=fill, font=font)
            return
        pad = max(tw, th) // 2 + 4
        stamp = Image.new("RGBA", (tw + pad * 2, th + pad * 2), (0, 0, 0, 0))
        ImageDraw.Draw(stamp).text((pad, pad), text, fill=fill, font=font)
        stamp = stamp.rotate(rotation, resample=Image.BICUBIC, expand=True)
        sw, sh = stamp.size
        cx, cy = x + tw // 2, y + th // 2
        layer.alpha_composite(stamp, (cx - sw // 2, cy - sh // 2))

    if position == "tiled":
        step_x = max(1, int((tw + 100) * max(0.3, tile_spacing)))
        step_y = max(1, int((th + 80) * max(0.3, tile_spacing)))
        for tx in range(0, w + step_x, step_x):
            for ty in range(0, h + step_y, step_y):
                _stamp(tx, ty)
    elif position == "custom":
        cx = int(custom_x * w) - tw // 2
        cy = int(custom_y * h) - th // 2
        cx = max(0, min(w - tw, cx))
        cy = max(0, min(h - th, cy))
        _stamp(cx, cy)
    elif position == "center":
        _stamp((w - tw) // 2, (h - th) // 2)
    elif position == "bottom-right":
        _stamp(w - tw - 20, h - th - 20)
    elif position == "bottom-left":
        _stamp(20, h - th - 20)
    elif position == "top-right":
        _stamp(w - tw - 20, 20)
    elif position == "top-left":
        _stamp(20, 20)
    else:
        _stamp((w - tw) // 2, (h - th) // 2)
    return Image.alpha_composite(img, layer)


def _scope_targets(scope, slide_num, slide_files):
    if scope == "current":
        try:
            n = int(slide_num)
        except (TypeError, ValueError):
            return None, "Invalid slide_num"
        if n < 1 or n > len(slide_files):
            return None, "Invalid slide_num"
        return [slide_files[n - 1]], None
    return list(slide_files), None


@app.route("/api/watermark", methods=["POST"])
def add_text_watermark():
    slide_files = _get_slide_files()
    if not slide_files:
        return jsonify({"error": "No slides loaded"}), 400
    payload = _ensure_dict(request.json)

    text = str(payload.get("text", "")).strip()[:200]
    if not text:
        return jsonify({"error": "Text cannot be empty"}), 400

    def _num(key, default, lo, hi):
        try:
            return max(lo, min(hi, float(payload.get(key, default))))
        except (TypeError, ValueError):
            return default
    opacity = _num("opacity", 0.25, 0.02, 1.0)
    font_scale = _num("font_scale", 1 / 14, 0.02, 0.25)
    rotation = _num("rotation", 0, -90, 90)
    tile_spacing = _num("tile_spacing", 1.0, 0.3, 4.0)
    custom_x = _num("custom_x", 0.5, 0.0, 1.0)
    custom_y = _num("custom_y", 0.5, 0.0, 1.0)
    position = str(payload.get("position", "center"))
    color = _parse_hex(payload.get("color", "#808080"))

    scope = str(payload.get("scope", "all"))
    targets, err = _scope_targets(scope, payload.get("slide_num"), slide_files)
    if err:
        return jsonify({"error": err}), 400

    snapshot = _snapshot("watermark-text")
    for sf in targets:
        img = Image.open(sf).convert("RGBA")
        img = _draw_text_watermark(
            img, text=text, color=color, opacity=opacity, position=position,
            font_scale=font_scale, rotation=rotation, tile_spacing=tile_spacing,
            custom_x=custom_x, custom_y=custom_y,
        )
        img.convert("RGB").save(str(sf), "JPEG", quality=95)

    log_id = _log_op("watermark-text",
                     text=f'Watermark "{text}" on {len(targets)} slide(s)',
                     snapshot=snapshot, count=len(targets))
    return jsonify({"ok": True, "count": len(targets),
                    "snapshot": snapshot, "log_id": log_id})


@app.route("/api/watermark-image", methods=["POST"])
def add_image_watermark():
    if "image" not in request.files:
        return jsonify({"error": "No image uploaded"}), 400
    try:
        wm_img = Image.open(request.files["image"].stream).convert("RGBA")
    except Exception:
        return jsonify({"error": "Invalid image"}), 400

    slide_files = _get_slide_files()
    if not slide_files:
        return jsonify({"error": "No slides loaded"}), 400

    def _num(key, default, lo, hi):
        try:
            return max(lo, min(hi, float(request.form.get(key, default))))
        except (TypeError, ValueError):
            return default
    opacity = _num("opacity", 0.3, 0.02, 1.0)
    scale = _num("scale", 0.15, 0.02, 0.6)
    custom_x = _num("custom_x", 0.5, 0.0, 1.0)
    custom_y = _num("custom_y", 0.5, 0.0, 1.0)
    position = request.form.get("position", "bottom-right")
    scope = request.form.get("scope", "all")

    try:
        slide_num = int(request.form.get("slide_num", 1))
    except (TypeError, ValueError):
        slide_num = 1
    targets, err = _scope_targets(scope, slide_num, slide_files)
    if err:
        return jsonify({"error": err}), 400

    snapshot = _snapshot("watermark-image")
    for sf in targets:
        img = Image.open(sf).convert("RGBA")
        w, h = img.size
        wm_w = int(w * scale)
        wm_h = max(1, int(wm_w * wm_img.height / max(1, wm_img.width)))
        wm_resized = wm_img.resize((wm_w, wm_h), Image.LANCZOS)
        r, g, b, a = wm_resized.split()
        a = a.point(lambda v: int(v * opacity))
        wm_resized = Image.merge("RGBA", (r, g, b, a))

        if position == "custom":
            cx = int(custom_x * w) - wm_w // 2
            cy = int(custom_y * h) - wm_h // 2
            cx = max(0, min(w - wm_w, cx))
            cy = max(0, min(h - wm_h, cy))
            pos = (cx, cy)
        elif position == "center":
            pos = ((w - wm_w) // 2, (h - wm_h) // 2)
        elif position == "top-left":
            pos = (20, 20)
        elif position == "top-right":
            pos = (w - wm_w - 20, 20)
        elif position == "bottom-left":
            pos = (20, h - wm_h - 20)
        elif position == "tiled":
            overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
            for tx in range(0, w, wm_w + 60):
                for ty in range(0, h, wm_h + 40):
                    overlay.paste(wm_resized, (tx, ty), wm_resized)
            img = Image.alpha_composite(img, overlay)
            img.convert("RGB").save(str(sf), "JPEG", quality=95)
            continue
        else:
            pos = (w - wm_w - 20, h - wm_h - 20)
        img.paste(wm_resized, pos, wm_resized)
        img.convert("RGB").save(str(sf), "JPEG", quality=95)

    log_id = _log_op("watermark-image",
                     text=f"Image watermark on {len(targets)} slide(s)",
                     snapshot=snapshot, count=len(targets))
    return jsonify({"ok": True, "count": len(targets),
                    "snapshot": snapshot, "log_id": log_id})


# ── Background cleanup of stale sessions ─────────────────────────────────
def _gc_sessions():
    cutoff = time.time() - SESSION_TTL_HOURS * 3600
    try:
        for d in SESSIONS_DIR.iterdir():
            if not d.is_dir():
                continue
            try:
                mtime = max((f.stat().st_mtime for f in d.rglob("*") if f.is_file()),
                            default=d.stat().st_mtime)
                if mtime < cutoff:
                    shutil.rmtree(str(d), ignore_errors=True)
            except OSError:
                continue
    except OSError:
        pass


def _gc_loop():
    while True:
        time.sleep(3600)
        _gc_sessions()


_gc_thread = threading.Thread(target=_gc_loop, daemon=True)
_gc_thread.start()


if __name__ == "__main__":
    host = os.environ.get('HOST', '127.0.0.1')
    port = int(os.environ.get('PORT', 5051))
    if host == '0.0.0.0':
        print("WARNING: binding to 0.0.0.0 — no auth.", file=sys.stderr)
    print(f"SlideCraft Lite on http://{host}:{port}", flush=True)
    app.run(debug=os.environ.get('FLASK_DEBUG', 'false').lower() == 'true',
            port=port, host=host)
